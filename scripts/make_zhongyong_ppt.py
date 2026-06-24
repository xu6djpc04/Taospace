# -*- coding: utf-8 -*-
"""
依中庸 NotebookLM 記事本內容建立簡報與設計稿文字
來源：115.04.24 概述首章.pptx / 114.11.02 中庸四[26~33].pptx / 2025-11-02 mp3
產出：
  output/中庸講座簡報_115年.pptx
  output/中庸講座設計稿_115年.md
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

OUTPUT_DIR = r"C:\Users\kevin\Desktop\Taospace\output"

DEEP_BROWN = RGBColor(0x4A, 0x2C, 0x00)
GOLD       = RGBColor(0xC8, 0xA0, 0x00)
CREAM      = RGBColor(0xFD, 0xF6, 0xE3)
DARK_TEXT  = RGBColor(0x2C, 0x1A, 0x00)
LIGHT_GOLD = RGBColor(0xF0, 0xD8, 0x80)

# 16:9 slide dimensions
W = 13.333
H = 7.5

SLIDES = [
    # 1 封面
    {
        "type": "cover",
        "title": "中 庸",
        "subtitle": "儒家修道之學　以誠立命之道",
    },
    # 2 何謂中庸
    {
        "title": "何謂「中庸」？",
        "subtitle": "中・庸 名稱解義",
        "body": [
            "中 = 本體（天下之大本）",
            "   不偏不倒，無過不及",
            "   喜怒哀樂未發——性體之謂",
            "庸 = 用（天下之達道）",
            "   率性而行，萬古不易之道",
            "   用中以明道，守中而不離道",
            "─────────────────",
            "中庸 = 依體起用・守中行道",
            "   「堆之則彌六合，巻之則退藏於密」",
        ],
    },
    # 3 朱子章句序
    {
        "title": "朱子章句序",
        "subtitle": "孔門傳授心法",
        "body": [
            "「不偏之謂中，不易之謂庸。」（程子）",
            "   中是天下人應循的正確道路",
            "   庸是天下事發展運行的永恆定律",
            "─────────────────",
            "「此書始言一理；中散為萬事；末復合為一理。」",
            "   放之則彌六合，巻之則退藏於密",
            "   其味無穷，皣實學也",
            "",
            "傳承脈絡：孔子 → 子思（筆之於書）→ 孟子",
        ],
    },
    # 4 首章三綱領
    {
        "title": "首章三綱領",
        "subtitle": "中庸第一章",
        "body": [
            "① 天命之謂性",
            "   上天賦予——靈明本性",
            "② 率性之謂道",
            "   依循天性直道而行",
            "③ 修道之謂教",
            "   品節本性，立教化人",
            "─────────────────",
            "「道也者，不可須臾離也；可離，非道也。」",
            "   道能保全人生、明人性，哪里可以片刻離開？",
        ],
    },
    # 5 修道入手工夫
    {
        "title": "修道入手工夫",
        "subtitle": "戒慎恐懼・慎其獨",
        "body": [
            "戒慎恐懼",
            "   在無人看見、無聲無響時仍警潔自持",
            "慎其獨",
            "   道在一念之微，修道貴修心，修心貴修念",
            "─────────────────",
            "「莫見乎隱，莫顯乎微，故君子慎其獨也。」",
            "",
            "起心動念皆記於宇宙記憶體",
            "→ 工夫在察覺念頭、降伏心魔",
        ],
    },
    # 6 致中和
    {
        "title": "致中和・天地位・萬物育",
        "subtitle": "中庸第一章（結語）",
        "body": [
            "中 = 喜怒哀樂未發（天下之大本）",
            "和 = 發而皆中節（天下之達道）",
            "─────────────────",
            "致中和（工夫）→ 天地位焉，萬物育焉（境界）",
            "",
            "首章三層次：",
            "   本體：天命之謂性",
            "   工夫：戒慎恐懼・慎其獨",
            "   境界：致中和→位天地・育萬物",
        ],
    },
    # 7 鬼神之為德
    {
        "title": "鬼神之為德",
        "subtitle": "中庸第十六章",
        "body": [
            "「鬼神之為德，其盛矣乎！」",
            "   此「鬼神」非一般所指，乃陰陽二氣",
            "   大至天地四季冷暖，小至自己內心隐微念頭",
            "─────────────────",
            "鬼神之別，在覺性與迷性之分",
            "   克念——神也（覺）",
            "   罔念——鬼也（迷）",
            "",
            "去鬼現神・捨姄顯真",
            "→ 一切存乎一心，成神成鬼皆由自己做主",
        ],
    },
    # 8 誠的核心思想
    {
        "title": "誠——中庸核心思想",
        "subtitle": "論語仁・中庸誠・體用兼含",
        "body": [
            "「誠」字兼含天道觀與人道思想",
            "   天道（本體）：誠者，天之道也",
            "   人道（工夫）：誠之者，人之道也",
            "─────────────────",
            "誠在天人存在意義上 → 本體、性體",
            "誠在人的表現上 → 化為諸德目",
            "",
            "誠故兼有體、用之義",
            "透過「誠」才能實現天人合一",
        ],
    },
    # 9 誠者天之道
    {
        "title": "誠者天之道　誠之者人之道",
        "subtitle": "第二十章・第二十一章",
        "body": [
            "誠者天之道",
            "   真實無姄，不勉而中，不思而得",
            "   從容中道——聖人（自誠明，謂之性）",
            "誠之者人之道",
            "   擇善而固執——賢人（自明誠，謂之教）",
            "─────────────────",
            "誠則明矣，明則誠矣",
            "   修道核心：實習修道方法",
            "               實踐戒慎之訓",
            "               實致中和之功",
        ],
    },
    # 10 誠心 vs 愚誠
    {
        "title": "誠心 vs 愕誠",
        "subtitle": "以智慧為基礎的誠",
        "body": [
            "誠心",
            "   有誠懇的心 + 有智慧的心",
            "   能看別人的想法，再看自己做的對不對",
            "愕誠",
            "   雖誠懇，卻依主觀意識、崇拜",
            "   視線窄，跟著別人對錯走",
            "─────────────────",
            "「信為道元功德母」",
            "能達到最後成就的關鍵因素——智慧",
            "   戒（守心）→ 定（定力）→ 慧（智慧）",
        ],
    },
    # 11 修誠三工夫
    {
        "title": "修誠三工夫",
        "subtitle": "擇善固執・致曲・知止",
        "body": [
            "擇善固執（第二十章）",
            "   依止自善寶地，守住不放失",
            "   博學・審問・慎思・明辨・篵行",
            "致曲（第二十三章）",
            "   委婉體察，有所受（容）・有所止（定）",
            "   止則不紛，不紛則一，一則靜，靜則安",
            "知止（大學工夫 / 中庸致曲之本）",
            "   停下來！止息姄心、過去牵掛、未來恐懼",
            "   不止難以得・難以定",
        ],
    },
    # 12 至誠如神
    {
        "title": "至誠如神",
        "subtitle": "第二十四章・感通天地",
        "body": [
            "「至誠之道，可以前知。」（第二十四章）",
            "   至誠者，通天命・知數理・窘人事",
            "─────────────────",
            "榮格求雨故事",
            "   老人將自身調整至與道合一，三日降雨",
            "師母至誠感通",
            "   每天一萬課手，關卡不解而解",
            "生公說法，頂石點頭",
            "   點道生的至誠感動天地萬物",
            "─────────────────",
            "「精誠所至，金石為開」",
        ],
    },
    # 13 至誠無息
    {
        "title": "至誠無息",
        "subtitle": "第二十六章",
        "body": [
            "「故至誠無息。」——真誠之道，永不止息",
            "─────────────────",
            "不息 → 久 → 徵 → 悝遠 → 博厘 → 高明",
            "   博厘配地，高明配天，悝久無疆",
            "",
            "「無息」入手工夫：調息攝心",
            "   氣不外泄，心不他移，然後可止可靜",
            "   「心一則神凝，神凝即至誠，至誠如神也」",
        ],
    },
    # 14 苟無至德
    {
        "title": "苟無至德　至道不凝",
        "subtitle": "第二十七～二十九章",
        "body": [
            "「大哉聖人之道！屸屸乎，發育萬物，峻極於天。」",
            "「苟無至德，至道不凝焉。」（第二十七章）",
            "   無至德，至道無法凝聚落實",
            "─────────────────",
            "三重：議禮・制度・考文",
            "   雖有其位，苟無其德，不敢作禮樂",
            "   雖有其德，苟無其位，亦不敢作禮樂",
            "─────────────────",
            "君子之道：本諸身，徵諸廫民",
            "   建諸天地而不悖，百世以俳聖人而不惑",
        ],
    },
    # 15 夫子至德配天
    {
        "title": "夫子至德　配天",
        "subtitle": "第三十～三十二章",
        "body": [
            "「仁尼祖述堯舜，憲章文武。」（第三十章）",
            "   上律天時，下襲水土",
            "   萬物並育而不相害，道並行而不相悖",
            "─────────────────",
            "「溥博淵泉，而時出之。」（第三十一章）",
            "   溥博如天；淵泉如淵",
            "   見而民莫不敬；言而民莫不信；行而民莫不說",
            "─────────────────",
            "「唯天下至誠，為能經綢天下之大經，",
            " 立天下之大本，知天地之化育。」（第三十二章）",
        ],
    },
    # 16 至精要義 / 總結
    {
        "title": "中庸之至精要義",
        "subtitle": "第三十三章（全篇結語）・返本歸誠",
        "body": [
            "「君子之道，闇然而日章；小人之道，的然而日亡。」",
            "「不顔惟德，百辟其刑之。」→ 君子篵恭而天下平",
            "─────────────────",
            "「上天之載，無聲無臭，至矣！」",
            "   道之極致——無形無聲，卻無所不在",
            "─────────────────",
            "中庸 = 守中行道，非平庸折中",
            "誠 = 天道本體，也是人道工夫",
            "至誠無息，化育天地——終極境界",
            "「道不遠人。人之為道而遠人，不可以為道。」",
        ],
    },
]


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=DARK_TEXT,
                align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "微軟正黑體"
    return txb


def add_body_textbox(slide, lines, left, top, width, height, font_size=17):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = DARK_TEXT
        run.font.name = "微軟正黑體"


def make_cover(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DEEP_BROWN)
    bar = slide.shapes.add_shape(1, Inches(0), Inches(3.0), Inches(W), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = GOLD; bar.line.fill.background()

    add_textbox(slide, data["title"],
                left=1.5, top=1.3, width=W - 3, height=1.5,
                font_size=60, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_textbox(slide, data["subtitle"],
                left=1.5, top=3.2, width=W - 3, height=0.8,
                font_size=22, color=CREAM, align=PP_ALIGN.CENTER)
    add_textbox(slide, "崇元培訓班　115年",
                left=0.5, top=6.9, width=W - 1, height=0.4,
                font_size=14, color=LIGHT_GOLD, align=PP_ALIGN.RIGHT)


def make_content(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CREAM)

    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(W), Inches(1.1))
    header.fill.solid(); header.fill.fore_color.rgb = DEEP_BROWN; header.line.fill.background()

    add_textbox(slide, data["title"],
                left=0.3, top=0.05, width=W - 0.6, height=0.9,
                font_size=28, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

    sub_bar = slide.shapes.add_shape(1, Inches(0), Inches(1.1), Inches(W), Inches(0.45))
    sub_bar.fill.solid(); sub_bar.fill.fore_color.rgb = LIGHT_GOLD; sub_bar.line.fill.background()

    add_textbox(slide, data.get("subtitle", ""),
                left=0.4, top=1.1, width=W - 0.8, height=0.45,
                font_size=16, color=DEEP_BROWN)

    if data.get("body"):
        add_body_textbox(slide, data["body"],
                         left=0.5, top=1.75, width=W - 1.0, height=5.1)

    footer = slide.shapes.add_shape(1, Inches(0), Inches(7.2), Inches(W), Inches(0.2))
    footer.fill.solid(); footer.fill.fore_color.rgb = DEEP_BROWN; footer.line.fill.background()


def build_pptx():
    prs = Presentation()
    prs.slide_width  = Inches(W)
    prs.slide_height = Inches(H)

    for data in SLIDES:
        if data.get("type") == "cover":
            make_cover(prs, data)
        else:
            make_content(prs, data)

    out_path = os.path.join(OUTPUT_DIR, "中庸講座簡報_115年.pptx")
    prs.save(out_path)
    print("OK pptx: " + out_path)


def build_md():
    text = """\
# 中庸講座簡報　設計稿

> 供 claude.ai/design 使用。16 張投影片，棕金配色，繁體中文。
> 配色：深棕 #4A2C00、金黃 #C8A000、米白 #FDF6E3

---

## Slide 1（封面）
**主標**：中　庸
**副標**：儒家修道之學　以誠立命之道
深棕背景，金色大字，右下角「崇元培訓班　115年」

---

## Slide 2　何謂「中庸」？
**副標**：中・庸 名稱解義

- **中** = 本體（天下之大本）—不偏不倒，喜怒哀樂未發之謂
- **庸** = 用（天下之達道）—率性而行，用中以明道
- 中庸 = 依體起用・守中行道
- 「深之則彌六合，巻之則退藏於密」

---

## Slide 3　朱子章句序
**副標**：孔門傳授心法

- 「不偏之謂中，不易之謂庸。」（程子）
- 「此書始言一理；中散為萬事；末復合為一理。」
- 放之則彌六合，巻之則退藏於密
- 傳承：孔子 → 子思 → 孟子

---

## Slide 4　首章三綱領
**副標**：中庸第一章

① 天命之謂性（上天賦予靈明本性）
② 率性之謂道（依循天性直道而行）
③ 修道之謂教（品節本性，立教化人）
「道也者，不可須臾離也；可離，非道也。」

---

## Slide 5　修道入手工夫
**副標**：戒慎恐懼・慎其獨

- **戒慎恐懼**：在無人看見、無聲無響時仍警潔自持
- **慎其獨**：道在一念之微，修道貴修心，修心貴修念
- 「莫見乎隱，莫顯乎微，故君子慎其獨也。」
- 起心動念皆記於宇宙記憶體 → 工夫在察覺念頭、降伏心魔

---

## Slide 6　致中和・天地位・萬物育
**副標**：中庸第一章（結語）

- **中** = 喜怒哀樂未發（天下之大本）
- **和** = 發而皆中節（天下之達道）
- 致中和（工夫）→ 天地位焉，萬物育焉（境界）
- 首章三層次：本體→工夫→境界

---

## Slide 7　鬼神之為德
**副標**：中庸第十六章

- 鬼神 = 陰陽二氣，非一般所指
- 克念——神也（覺）；罔念——鬼也（迷）
- 去鬼現神・捨姄顯真
- 一切存乎一心，成神成鬼皆由自己做主

---

## Slide 8　誠——中庸核心思想
**副標**：體用兼含

- 天道（本體）：誠者，天之道也
- 人道（工夫）：誠之者，人之道也
- 誠兼有體、用之義
- 透過誠，實現天人合一

---

## Slide 9　誠者天之道　誠之者人之道
**副標**：第二十・二十一章

- **誠者天之道**：不勉而中，不思而得——聖人（自誠明，謂之性）
- **誠之者人之道**：擇善而固執——賢人（自明誠，謂之教）
- 誠則明矣，明則誠矣
- 修道核心：實習修道方法 / 實踐戒慎之訓 / 實致中和之功

---

## Slide 10　誠心 vs 愕誠
**副標**：以智慧為基礎的誠

| | 誠心 | 愕誠 |
|--|--|--|
| 心 | 誠懇 + 智慧 | 誠懇但主觀 |
| 結果 | 看別人想法，再看自己 | 別人對錯都跟著走 |

能達到最後成就的關鍵因素——**智慧**
戒（守心）→ 定（定力）→ 慧（智慧）

---

## Slide 11　修誠三工夫
**副標**：擇善固執・致曲・知止

- **擇善固執**（第二十章）：依止自善寶地，博學・審問・慎思・明辨・篵行
- **致曲**（第二十三章）：委婉體察，有所受（容）・有所止（定）止則不紛
- **知止**：停下來！止息姄心、過去牵掛、未來恐懼
- 不止難以得・難以定

---

## Slide 12　至誠如神
**副標**：第二十四章・感通天地

- 「至誠之道，可以前知。」
- **榮格求雨**：老人將自身調整至與道合一，三日降雨
- **師母至誠**：每天一萬課手，關卡不解而解
- **生公說法，頂石點頭**：至誠感動天地萬物
- 「精誠所至，金石為開」

---

## Slide 13　至誠無息
**副標**：第二十六章

- 「故至誠無息。」——真誠之道，永不止息
- 不息 → 久 → 徵 → 悝遠 → 博厘 → 高明
- 博厘配地，高明配天，悝久無疆
- 「無息」工夫：調息攝心，氣不外泄，心不他移

---

## Slide 14　苟無至德　至道不凝
**副標**：第二十七～二十九章

- 「苟無至德，至道不凝焉。」
- 三重：議禮・制度・考文
- 德位相符，方能化民成俗
- 君子之道：本諸身，徵諸廫民

---

## Slide 15　夫子至德　配天
**副標**：第三十～三十二章

- 「仁尼祖述堯舜，憲章文武。」—上律天時，下襲水土
- 「溥博淵泉，而時出之。」—溥博如天；淵泉如淵
- 「唯天下至誠，為能經綢天下之大經」
- 萬物並育而不相害，道並行而不相悖

---

## Slide 16　中庸之至精要義（結語）
**副標**：第三十三章・返本歸誠

- 「君子之道，闇然而日章；小人之道，的然而日亡。」
- 「上天之載，無聲無臭，至矣！」—道之極致
- 中庸 = 守中行道，非平庸折中
- 誠 = 天道本體，也是人道工夫
- 至誠無息，化育天地——終極境界
- 「道不遠人。人之為道而遠人，不可以為道。」

---

## 設計建議
- **字型**：標題粗體無襯線（微軟正黑體 Bold），正文細體
- **配色**：深棕 `#4A2C00`（標題欄）、金黃 `#C8A000`（強調文字）、米白 `#FDF6E3`（背景）
- **版式**：頂部深棕横條放標題，主體米白，引句可加細線框或斌體區隔
- **引句排版**：建議加左側金色豪條作為 blockquote 樣式

---

*資料來源：中庸 NotebookLM 記事本（115年整理）*
"""
    out_path = os.path.join(OUTPUT_DIR, "中庸講座設計稿_115年.md")
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(text)
    print("OK md: " + out_path)


if __name__ == "__main__":
    build_pptx()
    build_md()
    print("done.")
