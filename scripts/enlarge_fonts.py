import unicodedata
from pptx import Presentation
from pptx.util import Pt
import os

def enlarge_fonts(in_path, out_path):
    prs = Presentation(unicodedata.normalize('NFC', in_path))
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    if run.font.size:
                        pt = run.font.size.pt
                        # 放大 4pt，但上限設為 32，避免字體過大跑版
                        if pt <= 25:
                            new_size = pt + 4
                            run.font.size = Pt(new_size)

    prs.save(unicodedata.normalize('NFC', out_path))

if __name__ == "__main__":
    in_file = r"C:\Users\kevin\Desktop\Taospace\output\中庸心得分享-總論精華版_115年.pptx"
    out_file = r"C:\Users\kevin\Desktop\Taospace\output\中庸心得分享-總論精華版_大字_115年.pptx"
    os.makedirs(os.path.dirname(out_file), exist_ok=True)
    enlarge_fonts(in_file, out_file)
    print("Done")
