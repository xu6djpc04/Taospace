import unicodedata
from pptx import Presentation
import os

def extract_text(pptx_path, out_file):
    path = unicodedata.normalize('NFC', pptx_path)
    try:
        prs = Presentation(path)
        out_file.write(f"--- {os.path.basename(pptx_path)} ({len(prs.slides)} slides) ---\n")
        for i, slide in enumerate(prs.slides):
            out_file.write(f"Slide {i+1}:\n")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    # replace newlines with spaces for compact view
                    text = shape.text.replace("\n", " ").replace("\r", " ")
                    if text.strip():
                        out_file.write(f"  {text}\n")
            out_file.write("\n")
    except Exception as e:
        out_file.write(f"Error reading {pptx_path}: {e}\n")

if __name__ == "__main__":
    os.makedirs(r"C:\Users\kevin\.gemini\antigravity-cli\brain\ba8c2f21-d47c-4bc4-ae0b-6d3d86de177a\scratch", exist_ok=True)
    out_path = unicodedata.normalize('NFC', r"C:\Users\kevin\.gemini\antigravity-cli\brain\ba8c2f21-d47c-4bc4-ae0b-6d3d86de177a\scratch\pptx_content.txt")
    
    with open(out_path, "w", encoding="utf-8-sig") as f:
        extract_text(r"C:\Users\kevin\Desktop\Taospace\Taofiles\中庸\中庸心得分享-劉坤翰.pptx", f)
        f.write("\n" + "="*50 + "\n\n")
        extract_text(r"C:\Users\kevin\Desktop\Taospace\Taofiles\中庸\中庸分享-劉坤翰.pptx", f)
        
    print(f"Extraction complete. Results saved to {out_path}")
