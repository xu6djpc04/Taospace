import unicodedata
from pptx import Presentation
import os

def delete_slide(prs, index):
    # index is 0-based
    slide_id = prs.slides._sldIdLst[index].rId
    prs.part.drop_rel(slide_id)
    del prs.slides._sldIdLst[index]

def main():
    source_path = unicodedata.normalize('NFC', r"C:\Users\kevin\Desktop\Taospace\Taofiles\中庸\中庸分享-劉坤翰.pptx")
    target_path = unicodedata.normalize('NFC', r"C:\Users\kevin\Desktop\Taospace\output\中庸心得分享-總論精華版_115年.pptx")
    
    # Ensure output directory exists
    os.makedirs(os.path.dirname(target_path), exist_ok=True)
    
    try:
        prs = Presentation(source_path)
        
        # Slides to KEEP (1-indexed): 1, 2, 3, 4, 5, 6, 8, 12, 16
        # Slides to DELETE (1-indexed): 7, 9, 10, 11, 13, 14, 15
        # Convert to 0-indexed for deletion: 6, 8, 9, 10, 12, 13, 14
        slides_to_delete = [6, 8, 9, 10, 12, 13, 14]
        
        # Sort in descending order to avoid index shifting when deleting
        slides_to_delete.sort(reverse=True)
        
        for idx in slides_to_delete:
            delete_slide(prs, idx)
            
        prs.save(target_path)
        print(f"SUCCESS: {target_path}")
    except Exception as e:
        print(f"ERROR: {e}")

if __name__ == "__main__":
    main()
