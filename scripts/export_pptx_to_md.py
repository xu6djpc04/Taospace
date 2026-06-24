# -*- coding: utf-8 -*-
"""
將「中庸心得分享-總論精華版_115年_v2.pptx」匯出為 Claude Design 設計簡報 MD
"""
import unicodedata
from pptx import Presentation
from pptx.util import Pt

SRC = unicodedata.normalize('NFC',
    r'C:\Users\kevin\Desktop\Taospace\output\中庸心得分享-總論精華版_115年_v2.pptx')
DST = unicodedata.normalize('NFC',
    r'C:\Users\kevin\Desktop\Taospace\output\中庸心得分享-總論精華版_設計簡報_115年.md')

prs = Presentation(SRC)

# 從原始設計稿推導的版型提示
LAYOUT_HINTS = {
    1: "封面版型：深棕背景，金色大字置中。右下角標註「崇元培訓班　115年」。",
    2: "左右雙欄：左欄放超大「中」「庸」字（裝飾用），右欄放定義說明。底部橫跨全寬放核心結論與引文。",
    3: "引文版型：上方兩段引文（朱子/程子），下方傳承圖（孔子→子思→孟子）橫排。",
    4: "三項縱列：① ② ③ 三個圓圈數字 + 綱領名稱 + 說明，底部引文橫跨全寬。",
    5: "兩段主體：「戒慎恐懼」與「慎其獨」各一段（粗體標題 + 說明），底部引文。",
    6: "三欄流程：「中」→「致中和（工夫）」→「和」橫排，底部「本體→工夫→境界」三層次說明。",
    7: "核心頁：背景有超大「誠」字（裝飾），前景左右雙欄（天道/人道），底部結論句。",
    8: "三卡片：頂部引文，下方三個故事卡（榮格求雨、師母至誠、生公說法），底部總結句。",
    9: "結語頁：背景有超大「誠」字（裝飾），前景三條要點 + 首尾引文。",
}

lines = []
lines.append("# 中庸心得分享－總論精華版　設計簡報")
lines.append("")
lines.append("> 此為給 Claude Design 的重新排版簡報。共 9 頁，約 10 分鐘分享用。")
lines.append("> 請依下列設計規格與各頁內容重新排版，**所有正文字體不得小於 35pt**。")
lines.append("")
lines.append("---")
lines.append("")
lines.append("## 設計規格")
lines.append("")
lines.append("| 項目 | 規格 |")
lines.append("|------|------|")
lines.append("| 配色 | 深棕 `#4A2C00`（主）、金黃 `#C8A000`（標題/強調）、米白 `#FDF6E3`（內文背景） |")
lines.append("| 字體 | 標題：Cambria 或 Georgia；內文：Calibri |")
lines.append("| 比例 | 16:9 寬螢幕 |")
lines.append("| 語言 | 繁體中文 |")
lines.append("| 最小字體 | **35pt**（裝飾性大字除外） |")
lines.append("| 風格 | 若單頁文字量過多，可拆成兩頁，以維持字體大小與留白 |")
lines.append("")
lines.append("---")
lines.append("")
lines.append("## 各頁內容")
lines.append("")

for slide_idx, slide in enumerate(prs.slides, 1):
    lines.append(f"### Slide {slide_idx}")
    lines.append("")

    # 版型提示
    hint = LAYOUT_HINTS.get(slide_idx, "")
    if hint:
        lines.append(f"**版型提示**：{hint}")
        lines.append("")

    # 文字內容（依 shape 順序，跳過超大裝飾字）
    lines.append("**內容**：")
    lines.append("")
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        block_texts = []
        for para in shape.text_frame.paragraphs:
            para_text = para.text.strip()
            if not para_text:
                continue
            # 判斷是否為裝飾大字（只有單一字且該 run size >= 60pt）
            is_deco = False
            for run in para.runs:
                if run.font.size and run.font.size.pt >= 60 and len(para_text) <= 2:
                    is_deco = True
            if is_deco:
                block_texts.append(f"〔裝飾大字〕{para_text}")
            else:
                block_texts.append(para_text)
        if block_texts:
            for t in block_texts:
                lines.append(f"- {t}")
            lines.append("")

    # 備注（講稿提示）
    if slide.has_notes_slide:
        note = slide.notes_slide.notes_text_frame.text.strip()
        if note:
            lines.append(f"**講師備注**：{note}")
            lines.append("")

    lines.append("---")
    lines.append("")

md_content = "\n".join(lines)

with open(DST, 'w', encoding='utf-8-sig') as f:
    f.write(md_content)

print(f"已輸出：{DST}")
