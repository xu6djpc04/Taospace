import unicodedata
from pptx import Presentation
import os

def inspect_fonts(pptx_path, out_file):
    path = unicodedata.normalize('NFC', pptx_path)
    prs = Presentation(path)
    with open(out_file, 'w', encoding='utf-8-sig') as f:
        for i, slide in enumerate(prs.slides):
            f.write(f"Slide {i+1}:\n")
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        size = run.font.size.pt if run.font.size else "None (Inherited)"
                        text = run.text.strip()
                        if text:
                            f.write(f"  [Size: {size}] {text}\n")

if __name__ == "__main__":
    out_path = r"C:\Users\kevin\.gemini\antigravity-cli\brain\ba8c2f21-d47c-4bc4-ae0b-6d3d86de177a\scratch\font_sizes.txt"
    inspect_fonts(r"C:\Users\kevin\Desktop\Taospace\output\中庸心得分享-總論精華版_115年.pptx", out_path)
    print("Done")
