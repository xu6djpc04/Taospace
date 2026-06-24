# -*- coding: utf-8 -*-
"""
修正「中庸心得分享-總論精華版」PPTX：
1. 修正 Slide 1 備注「本講共十六講」→ 刪除，避免與 9 頁版本矛盾
2. 字體最小值提升：< 20pt → 20pt；20~23pt → 24pt（大裝飾字不動）
"""
import unicodedata
from pptx import Presentation
from pptx.util import Pt

SRC = unicodedata.normalize('NFC', r'C:\Users\kevin\Desktop\Taospace\output\中庸心得分享-總論精華版_115年.pptx')
DST = unicodedata.normalize('NFC', r'C:\Users\kevin\Desktop\Taospace\output\中庸心得分享-總論精華版_115年_v2.pptx')

prs = Presentation(SRC)

changed = []

for slide_idx, slide in enumerate(prs.slides, 1):
    # --- 1. 修正 Slide 1 備注 ---
    if slide_idx == 1 and slide.has_notes_slide:
        tf = slide.notes_slide.notes_text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                if '本講共十六講，' in run.text:
                    run.text = run.text.replace('本講共十六講，', '')
                    changed.append(f'Slide 1 備注：移除「本講共十六講，」')

    # --- 2. 字體大小調整 ---
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size is None:
                    continue
                pt = run.font.size.pt
                if pt >= 40:   # 裝飾大字（中、庸、誠等）不動
                    continue
                if pt < 20:
                    run.font.size = Pt(20)
                    changed.append(f'Slide {slide_idx}: {pt}pt → 20pt | {run.text[:30]}')
                elif pt < 24:
                    run.font.size = Pt(24)
                    changed.append(f'Slide {slide_idx}: {pt}pt → 24pt | {run.text[:30]}')

prs.save(DST)

print('=== 修改清單 ===')
for c in changed:
    print(' ', c)
print(f'\n已儲存：{DST}')
